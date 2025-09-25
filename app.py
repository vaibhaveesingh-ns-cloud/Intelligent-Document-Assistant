import os
import io
import re
import uuid
import hashlib
from dataclasses import dataclass
from datetime import datetime
from typing import Any, Dict, List, Optional, Set, Tuple

# Load environment variables from .env file
from dotenv import load_dotenv
load_dotenv()

import streamlit as st

st.set_page_config(page_title="Intelligent Document Assistant", page_icon="🧠", layout="wide")

import pandas as pd
from PIL import Image

# Optional dependencies handled gracefully
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pytesseract
except Exception:
    pytesseract = None

# LangChain imports
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document
    from langchain_openai import OpenAIEmbeddings, ChatOpenAI
    from langchain_chroma import Chroma
    from langchain.chains import ConversationalRetrievalChain
    from langchain.memory import ConversationBufferWindowMemory
    from langchain.prompts import PromptTemplate
except Exception as e:
    st.error(f"LangChain dependencies not installed: {e}")
    st.stop()

# Chroma vector database
try:
    import chromadb
except Exception:
    st.error("ChromaDB not installed. Run: pip install chromadb")
    st.stop()

# -----------------------------
# Utilities and Configuration
# -----------------------------

def clean_text(text: str) -> str:
    """Clean and normalize text content."""
    text = re.sub(r"\s+", " ", text)
    return text.strip()


@dataclass
class ConversationTurn:
    """Represents a single conversation turn."""
    question: str
    answer: str
    sources: List[str]
    timestamp: datetime


class DocumentProcessor:
    """Enhanced document processing with LangChain integration."""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

    def process_documents(self, texts_and_metadata: List[Tuple[str, Dict[str, Any]]]) -> List[Document]:
        """Process texts into LangChain documents with metadata."""
        documents = []
        for text, metadata in texts_and_metadata:
            if not text.strip():
                continue

            # Create document with metadata
            doc = Document(
                page_content=clean_text(text),
                metadata={
                    **metadata,
                    'processed_at': datetime.now().isoformat(),
                    'doc_id': str(uuid.uuid4())
                }
            )

            # Split into chunks
            chunks = self.text_splitter.split_documents([doc])

            # Add chunk-specific metadata
            for i, chunk in enumerate(chunks):
                chunk.metadata.update({
                    'chunk_id': f"{chunk.metadata['doc_id']}_chunk_{i}",
                    'chunk_index': i,
                    'total_chunks': len(chunks)
                })

            documents.extend(chunks)

        return documents


# -----------------------------
# Loaders for different file types
# -----------------------------

def load_pdf(file: io.BytesIO) -> str:
    if pdfplumber is None:
        st.warning("pdfplumber not installed. Run: pip install pdfplumber")
        return ""
    all_text = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            all_text.append(page.extract_text() or "")
    return "\n".join(all_text)


def load_docx(file: io.BytesIO) -> str:
    if docx is None:
        st.warning("python-docx not installed. Run: pip install python-docx")
        return ""
    d = docx.Document(file)
    return "\n".join([p.text for p in d.paragraphs])


def load_pptx(file: io.BytesIO) -> str:
    if Presentation is None:
        st.warning("python-pptx not installed. Run: pip install python-pptx")
        return ""
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def load_image(file: io.BytesIO) -> str:
    if pytesseract is None:
        st.warning("pytesseract not installed (and requires Tesseract binary). Skipping OCR.")
        return ""
    try:
        image = Image.open(file)
        return pytesseract.image_to_string(image)
    except Exception:
        return ""


def load_txt(file: io.BytesIO) -> str:
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception:
        return file.read().decode("latin-1", errors="ignore")


def load_md(file: io.BytesIO) -> str:
    return load_txt(file)


def load_csv(file: io.BytesIO) -> str:
    try:
        df = pd.read_csv(file)
    except Exception:
        file.seek(0)
        df = pd.read_excel(file)
    # Convert to a readable text table
    return df.to_csv(index=False)


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
    ".md": load_md,
    ".csv": load_csv,
    ".xlsx": load_csv,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
}


# -----------------------------
# Enhanced Vector Store with Chroma
# -----------------------------


class EnhancedVectorStore:
    """Thin wrapper around a persistent Chroma vector store."""

    def __init__(
        self,
        collection_name: str,
        embedding_model: str,
        use_openai_embeddings: bool = False,
        persist_path: str = "./chroma_db",
    ) -> None:
        self.collection_name = collection_name
        self.embedding_model = embedding_model
        self.use_openai_embeddings = use_openai_embeddings
        self.persist_path = persist_path
        self.embeddings = self._initialize_embeddings()
        self.client = chromadb.PersistentClient(path=self.persist_path)
        self.vectorstore = self._create_vectorstore()

    def _initialize_embeddings(self):
        """Create an embedding function for the vector store. Uses OpenAI only to avoid heavy local models."""

        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            st.error("OpenAI API key is required for embeddings. Set OPENAI_API_KEY in your .env file.")
            st.stop()

        try:
            return OpenAIEmbeddings(model=os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-small"))
        except Exception as exc:  # pragma: no cover - runtime safety
            st.error(f"Failed to initialize OpenAI embeddings: {exc}")
            st.stop()

    def _create_vectorstore(self):
        """Create or retrieve a persistent Chroma collection."""

        try:
            return Chroma(
                client=self.client,
                collection_name=self.collection_name,
                embedding_function=self.embeddings,
            )
        except Exception as exc:  # pragma: no cover - hard failure
            st.error(f"Failed to initialize Chroma vector store: {exc}")
            st.stop()

    def add_documents(self, documents: List[Document]) -> int:
        """Add documents to the vector store and return how many were indexed."""

        if not documents:
            return 0

        try:
            self.vectorstore.add_documents(documents)
            return len(documents)
        except Exception as exc:
            st.error(f"Failed to add documents to the vector store: {exc}")
            return 0

    def similarity_search_with_score(self, query: str, k: int = 5):
        """Retrieve the top-k documents with similarity scores."""

        try:
            return self.vectorstore.similarity_search_with_score(query, k=k)
        except Exception as exc:
            st.error(f"Vector search failed: {exc}")
            return []

    def get_retriever(self, search_kwargs: Optional[Dict[str, Any]] = None):
        """Expose the LangChain retriever interface."""

        if search_kwargs is None:
            search_kwargs = {"k": 5}
        return self.vectorstore.as_retriever(search_kwargs=search_kwargs)

    def get_collection_count(self) -> int:
        """Return the number of vectors stored in the underlying collection."""

        try:
            return self.vectorstore._collection.count()
        except Exception:
            return 0

    def clear(self) -> None:
        """Remove the current collection and recreate an empty one."""

        try:
            self.client.delete_collection(self.collection_name)
        except Exception:
            # If the collection does not exist we can safely ignore the error.
            pass
        self.vectorstore = self._create_vectorstore()


# -----------------------------
# Vector index that handles chunking & ingestion
# -----------------------------


class VectorIndex:
    """Helper around the enhanced vector store to manage document ingestion."""

    def __init__(
        self,
        collection_name: str,
        model_name: str,
        chunk_size: int,
        chunk_overlap: int,
        use_openai_embeddings: bool = False,
    ) -> None:
        self.collection_name = collection_name
        self.model_name = model_name
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.vector_store = EnhancedVectorStore(
            collection_name=collection_name,
            embedding_model=model_name,
            use_openai_embeddings=use_openai_embeddings,
        )
        self.document_processor = DocumentProcessor(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )

    def update_chunking(self, chunk_size: int, chunk_overlap: int) -> None:
        """Update the chunking configuration for future ingestions."""

        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.document_processor = DocumentProcessor(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )

    def add_documents(self, texts_and_metadata: List[Tuple[str, Dict[str, Any]]]) -> Tuple[int, Dict[str, int]]:
        """Add raw texts with metadata and return chunk statistics."""

        documents = self.document_processor.process_documents(texts_and_metadata)
        indexed = self.vector_store.add_documents(documents)
        per_source: Dict[str, int] = {}

        for doc in documents:
            source_name = doc.metadata.get("source", "Unknown source")
            per_source[source_name] = per_source.get(source_name, 0) + 1

        return indexed, per_source

    def similarity_search_with_score(self, query: str, k: int = 5):
        return self.vector_store.similarity_search_with_score(query, k=k)

    def get_retriever(self, k: int = 5):
        return self.vector_store.get_retriever({"k": k})

    def get_collection_count(self) -> int:
        return self.vector_store.get_collection_count()


# -----------------------------
# Conversational QA system with memory & citations
# -----------------------------


class ConversationalQASystem:
    """Conversational question answering on top of the indexed documents."""

    def __init__(self, vector_index: VectorIndex, top_k: int = 5, memory_window: int = 6) -> None:
        self.vector_index = vector_index
        self.top_k = top_k
        self.memory_window = memory_window
        self.conversation_memory = ConversationBufferWindowMemory(
            k=memory_window,
            memory_key="chat_history",
            input_key="question",
            output_key="answer",
            return_messages=True,
        )
        self.qa_chain = None
        self._initialize_qa_chain()

    def _initialize_qa_chain(self) -> None:
        api_key = os.getenv("OPENAI_API_KEY")

        if not api_key:
            self.qa_chain = None
            return

        try:
            llm = ChatOpenAI(
                model_name=os.getenv("OPENAI_MODEL", "gpt-4o-mini"),
                temperature=0.2,
                api_key=api_key,
            )

            custom_prompt = PromptTemplate(
                template="""
You are an intelligent document assistant. Answer the question using ONLY the provided context.

Requirements:
- Ground every statement in the retrieved context.
- Cite each supporting document inline using the format [Source: <filename>].
- If multiple documents contribute, cite each of them.
- If the answer is not in the context, say that the documents do not contain the information.

Conversation so far:
{chat_history}

Context:
{context}

Question: {question}

Answer:
""".strip(),
                input_variables=["context", "chat_history", "question"],
            )

            retriever = self.vector_index.get_retriever(self.top_k)

            self.qa_chain = ConversationalRetrievalChain.from_llm(
                llm=llm,
                retriever=retriever,
                combine_docs_chain_kwargs={"prompt": custom_prompt},
                return_source_documents=True,
            )
        except Exception as exc:  # pragma: no cover - runtime safety
            st.error(f"Failed to initialize the OpenAI chat model: {exc}")
            self.qa_chain = None

    def reset_memory(self) -> None:
        self.conversation_memory.clear()

    def get_answer(self, question: str) -> Dict[str, Any]:
        if not self.qa_chain:
            return {
                "answer": "Please provide an OpenAI API key to enable answer generation.",
                "source_documents": [],
            }

        chat_history = self.conversation_memory.load_memory_variables({}).get("chat_history", [])

        try:
            if hasattr(self.qa_chain, "invoke"):
                result = self.qa_chain.invoke({"question": question, "chat_history": chat_history})
            else:  # Fallback for older LangChain versions
                result = self.qa_chain({"question": question, "chat_history": chat_history})
        except Exception as exc:
            return {
                "answer": f"Error generating answer: {exc}",
                "source_documents": [],
            }

        self.conversation_memory.save_context({"question": question}, {"answer": result.get("answer", "")})
        return result

    @staticmethod
    def build_citations(source_documents: List[Document]) -> List[str]:
        citations: List[str] = []
        seen: Set[str] = set()

        for doc in source_documents or []:
            metadata = doc.metadata or {}
            source_name = metadata.get("source", "Unknown source")
            page = metadata.get("page")
            chunk_index = metadata.get("chunk_index")
            total_chunks = metadata.get("total_chunks")

            details = source_name
            if page is not None:
                details += f", page {page}"
            if chunk_index is not None and total_chunks is not None:
                details += f", chunk {chunk_index + 1}/{total_chunks}"

            citation = f"[Source: {details}]"
            if citation not in seen:
                seen.add(citation)
                citations.append(citation)

        return citations


# -----------------------------
# Helper utilities for Streamlit state management
# -----------------------------


def ensure_state_defaults() -> None:
    st.session_state.setdefault("conversation_turns", [])
    st.session_state.setdefault("last_contexts", [])
    st.session_state.setdefault("indexed_file_hashes", set())


def ensure_vector_index(model_name: str, chunk_size: int, chunk_overlap: int) -> VectorIndex:
    ensure_state_defaults()

    settings = {
        "model_name": model_name,
        "chunk_size": chunk_size,
        "chunk_overlap": chunk_overlap,
    }

    reset_index = False

    if "index_settings" not in st.session_state:
        reset_index = True
    elif st.session_state.get("index_settings") != settings:
        reset_index = True

    if reset_index or "vector_index" not in st.session_state:
        st.session_state["index_settings"] = settings
        st.session_state["collection_name"] = f"documents_{uuid.uuid4().hex[:8]}"
        st.session_state["vector_index"] = VectorIndex(
            collection_name=st.session_state["collection_name"],
            model_name=model_name,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            use_openai_embeddings=True,
        )
        st.session_state["conversation_turns"] = []
        st.session_state["last_contexts"] = []
        st.session_state["indexed_file_hashes"] = set()
        st.session_state["qa_system"] = None

    return st.session_state["vector_index"]


def ensure_qa_system(vector_index: VectorIndex, top_k: int, memory_window: int) -> ConversationalQASystem:
    qa_system: Optional[ConversationalQASystem] = st.session_state.get("qa_system")

    api_key_available = bool(os.getenv("OPENAI_API_KEY"))
    if st.session_state.get("api_key_available") != api_key_available:
        st.session_state["api_key_available"] = api_key_available
        qa_system = None

    if (
        qa_system is None
        or qa_system.vector_index is not vector_index
        or qa_system.top_k != top_k
        or qa_system.memory_window != memory_window
    ):
        qa_system = ConversationalQASystem(vector_index, top_k=top_k, memory_window=memory_window)
        st.session_state["qa_system"] = qa_system

    return qa_system


def ingest_uploaded_files(uploaded_files, vector_index: VectorIndex):
    texts_and_metadata: List[Tuple[str, Dict[str, Any]]] = []
    skipped_duplicates: List[str] = []
    unsupported_files: List[str] = []
    empty_files: List[str] = []
    failed_files: List[Tuple[str, str]] = []
    ingested_hashes: Set[str] = st.session_state.setdefault("indexed_file_hashes", set())
    pending_hashes: List[str] = []

    for uploaded_file in uploaded_files:
        filename = uploaded_file.name
        suffix = os.path.splitext(filename)[1].lower()
        loader = LOADER_MAP.get(suffix)

        if not loader:
            unsupported_files.append(filename)
            continue

        file_bytes = uploaded_file.getvalue()
        file_hash = hashlib.md5(file_bytes).hexdigest()

        if file_hash in ingested_hashes:
            skipped_duplicates.append(filename)
            continue

        buffer = io.BytesIO(file_bytes)
        buffer.seek(0)

        try:
            raw_text = loader(buffer)
        except Exception as exc:
            failed_files.append((filename, str(exc)))
            continue

        cleaned_text = clean_text(raw_text)
        if not cleaned_text:
            empty_files.append(filename)
            continue

        metadata = {
            "source": filename,
            "file_type": suffix.lstrip("."),
            "file_hash": file_hash,
            "ingested_at": datetime.utcnow().isoformat(),
            "chunk_size": vector_index.chunk_size,
            "chunk_overlap": vector_index.chunk_overlap,
        }

        texts_and_metadata.append((cleaned_text, metadata))
        pending_hashes.append(file_hash)

    if not texts_and_metadata:
        return {
            "total_chunks": 0,
            "per_source": {},
            "skipped": skipped_duplicates,
            "unsupported": unsupported_files,
            "empty": empty_files,
            "failed": failed_files,
        }

    total_chunks, per_source = vector_index.add_documents(texts_and_metadata)

    if total_chunks:
        for file_hash in pending_hashes:
            ingested_hashes.add(file_hash)

    return {
        "total_chunks": total_chunks,
        "per_source": per_source,
        "skipped": skipped_duplicates,
        "unsupported": unsupported_files,
        "empty": empty_files,
        "failed": failed_files,
    }


def build_context_summaries(similarity_results) -> List[Dict[str, Any]]:
    summaries: List[Dict[str, Any]] = []

    for item in similarity_results:
        if isinstance(item, tuple) and len(item) == 2:
            doc, score = item
        else:
            doc, score = item, None

        metadata = doc.metadata or {}
        summaries.append(
            {
                "source": metadata.get("source", "Unknown source"),
                "chunk_index": metadata.get("chunk_index"),
                "total_chunks": metadata.get("total_chunks"),
                "score": score,
                "text": doc.page_content,
            }
        )

    return summaries


# -----------------------------
# Streamlit interface
# -----------------------------


st.title("🧠 Intelligent Document Assistant")
st.caption(
    "Upload multi-format documents, index them into a vector database, and ask grounded questions with cited answers."
)

with st.sidebar:
    st.header("Configuration")
    
    # Check if API key is available from environment
    api_key_available = bool(os.getenv("OPENAI_API_KEY"))
    if api_key_available:
        st.success("✅ OpenAI API key loaded from environment")
        st.info("Using OpenAI embeddings (no PyTorch required)")
    else:
        st.error("❌ OpenAI API key not found in .env file")
        st.markdown("""
        **To enable answer generation:**
        1. Copy `.env.example` to `.env`
        2. Add your OpenAI API key to the `.env` file
        3. Restart the application
        """)

    # Embedding model name for OpenAI (optional override via env)
    model_name = os.getenv("OPENAI_EMBED_MODEL", "text-embedding-3-small")

    chunk_size = st.slider("Chunk size (characters)", min_value=500, max_value=3000, value=1200, step=100)
    chunk_overlap = st.slider("Chunk overlap", min_value=0, max_value=600, value=200, step=50)
    if chunk_overlap >= chunk_size:
        st.info("Chunk overlap adjusted to be smaller than the chunk size.")
        chunk_overlap = max(0, chunk_size - 200)

    top_k = st.slider("Top-K chunks to retrieve", min_value=1, max_value=10, value=5)
    memory_window = st.slider(
        "Conversation memory (turns)",
        min_value=1,
        max_value=20,
        value=6,
        help="How many previous exchanges to keep in the conversation buffer.",
    )

    st.markdown("---")
    if st.button("Clear chat history"):
        ensure_state_defaults()
        st.session_state["conversation_turns"] = []
        st.session_state["last_contexts"] = []
        qa_system = st.session_state.get("qa_system")
        if qa_system:
            qa_system.reset_memory()
        st.success("Conversation cleared.")

    if st.button("Reset indexed documents"):
        for key in ["vector_index", "index_settings", "collection_name"]:
            st.session_state.pop(key, None)
        st.session_state["indexed_file_hashes"] = set()
        st.session_state["conversation_turns"] = []
        st.session_state["last_contexts"] = []
        st.session_state["qa_system"] = None
        st.success("Vector index reset. Upload documents again to re-index.")

vector_index = ensure_vector_index(model_name, chunk_size, chunk_overlap)
qa_system = ensure_qa_system(vector_index, top_k, memory_window)

col_main, col_info = st.columns([2, 1])

with col_main:
    st.subheader("1. Upload & index your documents")
    uploaded_files = st.file_uploader(
        "Upload documents",
        type=list(LOADER_MAP.keys()),
        accept_multiple_files=True,
        help="PDF, Word, PowerPoint, text, markdown, spreadsheets, and common image formats are supported.",
    )

    if uploaded_files:
        st.info("Click **Index documents** after selecting files to add them to the knowledge base.")
        if st.button("Index documents", use_container_width=True):
            with st.spinner("Extracting text and updating the vector database..."):
                ingestion_result = ingest_uploaded_files(uploaded_files, vector_index)

            total_chunks = ingestion_result["total_chunks"]
            if total_chunks:
                st.success(
                    f"Indexed {total_chunks} chunk(s) from {len(ingestion_result['per_source'])} document(s)."
                )
                for source_name, count in ingestion_result["per_source"].items():
                    st.caption(f"• {source_name}: {count} chunk(s)")
            else:
                st.info("No new content was indexed.")

            if ingestion_result["unsupported"]:
                st.warning(
                    "Unsupported file types: " + ", ".join(sorted(set(ingestion_result["unsupported"])))
                )
            if ingestion_result["skipped"]:
                st.info("Skipped previously indexed files: " + ", ".join(ingestion_result["skipped"]))
            if ingestion_result["empty"]:
                st.warning("No text detected in: " + ", ".join(ingestion_result["empty"]))
            if ingestion_result["failed"]:
                for filename, error_msg in ingestion_result["failed"]:
                    st.error(f"Failed to process {filename}: {error_msg}")

    st.divider()
    st.subheader("2. Chat with your documents")
    chat_container = st.container()
    context_container = st.container()

with chat_container:
    if not st.session_state["conversation_turns"]:
        st.info("Ask your first question to begin the conversation.")
    else:
        for turn in st.session_state["conversation_turns"]:
            with st.chat_message("user"):
                st.markdown(turn.question)
            with st.chat_message("assistant"):
                st.markdown(turn.answer)
                if turn.sources:
                    st.markdown("**Sources:** " + ", ".join(turn.sources))
                st.caption(turn.timestamp.strftime("%Y-%m-%d %H:%M:%S"))

with context_container:
    if st.session_state["last_contexts"]:
        st.markdown("### Retrieved context from the last answer")
        for idx, ctx in enumerate(st.session_state["last_contexts"], start=1):
            title = f"[{idx}] {ctx['source']}"
            if ctx.get("score") is not None:
                title += f" — similarity {ctx['score']:.3f}"
            with st.expander(title, expanded=False):
                meta_bits = []
                if ctx.get("chunk_index") is not None and ctx.get("total_chunks") is not None:
                    meta_bits.append(f"Chunk {ctx['chunk_index'] + 1} of {ctx['total_chunks']}")
                if meta_bits:
                    st.caption(" | ".join(meta_bits))
                st.write(ctx["text"])

user_question = st.chat_input("Ask a question about your indexed documents")

if user_question:
    with chat_container:
        with st.chat_message("user"):
            st.markdown(user_question)

    if vector_index.get_collection_count() == 0:
        assistant_response = "Please upload and index documents before asking questions."
        with chat_container:
            with st.chat_message("assistant"):
                st.markdown(assistant_response)
        st.session_state["conversation_turns"].append(
            ConversationTurn(
                question=user_question,
                answer=assistant_response,
                sources=[],
                timestamp=datetime.utcnow(),
            )
        )
    else:
        result = qa_system.get_answer(user_question)
        answer_text = result.get("answer", "I wasn't able to generate an answer.")
        citations = ConversationalQASystem.build_citations(result.get("source_documents", []))

        with chat_container:
            with st.chat_message("assistant"):
                st.markdown(answer_text)
                if citations:
                    st.markdown("**Sources:** " + ", ".join(citations))

        st.session_state["conversation_turns"].append(
            ConversationTurn(
                question=user_question,
                answer=answer_text,
                sources=citations,
                timestamp=datetime.utcnow(),
            )
        )

        similarity_results = vector_index.similarity_search_with_score(user_question, k=top_k)
        st.session_state["last_contexts"] = build_context_summaries(similarity_results)

        with context_container:
            context_container.empty()
            if st.session_state["last_contexts"]:
                st.markdown("### Retrieved context from the last answer")
                for idx, ctx in enumerate(st.session_state["last_contexts"], start=1):
                    title = f"[{idx}] {ctx['source']}"
                    if ctx.get("score") is not None:
                        title += f" — similarity {ctx['score']:.3f}"
                    with st.expander(title, expanded=False):
                        meta_bits = []
                        if ctx.get("chunk_index") is not None and ctx.get("total_chunks") is not None:
                            meta_bits.append(
                                f"Chunk {ctx['chunk_index'] + 1} of {ctx['total_chunks']}"
                            )
                        if meta_bits:
                            st.caption(" | ".join(meta_bits))
                        st.write(ctx["text"])

with col_info:
    st.subheader("Index status")
    st.metric("Chunks indexed", vector_index.get_collection_count())
    st.caption(
        "Indexing is session-scoped. Changing the embedding model or chunking parameters rebuilds the collection."
    )

    st.markdown("---")
    st.subheader("Conversation")
    st.metric("Stored turns", len(st.session_state["conversation_turns"]))
    st.caption("Conversation history is kept locally in your browser session and is never uploaded.")

    st.markdown("---")
    st.subheader("Tips")
    st.markdown(
        """
- Provide an **OpenAI API key** to enable grounded answer generation. Without it you'll still see retrieved context.
- Adjust **chunk size/overlap** based on the document type. Smaller chunks work well for slides, larger for reports.
- Image ingestion uses OCR via `pytesseract`; ensure the Tesseract binary is installed locally if you plan to use it.
- Chroma persists to `./chroma_db`. Reset the index from the sidebar to start fresh.
        """
    )
