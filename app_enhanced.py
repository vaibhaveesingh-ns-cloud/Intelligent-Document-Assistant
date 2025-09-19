import os
import io
import re
import time
import uuid
from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional
from datetime import datetime

import streamlit as st

# Parsing
import pandas as pd
from PIL import Image

# Optional dependencies handled gracefully
try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pytesseract
except Exception:
    pytesseract = None

# LangChain imports
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.schema import Document
    from langchain_openai import OpenAIEmbeddings, ChatOpenAI
    from langchain_chroma import Chroma
    from langchain.chains import ConversationalRetrievalChain
    from langchain.memory import ConversationBufferWindowMemory
    from langchain.prompts import PromptTemplate
    from langchain_community.embeddings import SentenceTransformerEmbeddings
except Exception as e:
    st.error(f"LangChain dependencies not installed: {e}")
    st.stop()

# Chroma vector database
try:
    import chromadb
    from chromadb.config import Settings
except Exception:
    st.error("ChromaDB not installed. Run: pip install chromadb")
    st.stop()

# OpenAI
try:
    from openai import OpenAI
except Exception:
    OpenAI = None

# -----------------------------
# Utilities and Configuration
# -----------------------------

def clean_text(text: str) -> str:
    """Clean and normalize text content."""
    text = re.sub(r"\s+", " ", text)
    return text.strip()


@dataclass
class ConversationTurn:
    """Represents a single conversation turn."""
    question: str
    answer: str
    sources: List[str]
    timestamp: datetime


class DocumentProcessor:
    """Enhanced document processing with LangChain integration."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def process_documents(self, texts_and_metadata: List[Tuple[str, Dict]]) -> List[Document]:
        """Process texts into LangChain documents with metadata."""
        documents = []
        for text, metadata in texts_and_metadata:
            if not text.strip():
                continue
            
            # Create document with metadata
            doc = Document(
                page_content=clean_text(text),
                metadata={
                    **metadata,
                    'processed_at': datetime.now().isoformat(),
                    'doc_id': str(uuid.uuid4())
                }
            )
            
            # Split into chunks
            chunks = self.text_splitter.split_documents([doc])
            
            # Add chunk-specific metadata
            for i, chunk in enumerate(chunks):
                chunk.metadata.update({
                    'chunk_id': f"{chunk.metadata['doc_id']}_chunk_{i}",
                    'chunk_index': i,
                    'total_chunks': len(chunks)
                })
            
            documents.extend(chunks)
        
        return documents


# -----------------------------
# Loaders for different file types
# -----------------------------

def load_pdf(file: io.BytesIO) -> str:
    if pdfplumber is None:
        st.warning("pdfplumber not installed. Run: pip install pdfplumber")
        return ""
    all_text = []
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            all_text.append(page.extract_text() or "")
    return "\n".join(all_text)


def load_docx(file: io.BytesIO) -> str:
    if docx is None:
        st.warning("python-docx not installed. Run: pip install python-docx")
        return ""
    d = docx.Document(file)
    return "\n".join([p.text for p in d.paragraphs])


def load_pptx(file: io.BytesIO) -> str:
    if Presentation is None:
        st.warning("python-pptx not installed. Run: pip install python-pptx")
        return ""
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def load_image(file: io.BytesIO) -> str:
    if pytesseract is None:
        st.warning("pytesseract not installed (and requires Tesseract binary). Skipping OCR.")
        return ""
    try:
        image = Image.open(file)
        return pytesseract.image_to_string(image)
    except Exception:
        return ""


def load_txt(file: io.BytesIO) -> str:
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception:
        return file.read().decode("latin-1", errors="ignore")


def load_md(file: io.BytesIO) -> str:
    return load_txt(file)


def load_csv(file: io.BytesIO) -> str:
    try:
        df = pd.read_csv(file)
    except Exception:
        file.seek(0)
        df = pd.read_excel(file)
    # Convert to a readable text table
    return df.to_csv(index=False)


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
    ".md": load_md,
    ".csv": load_csv,
    ".xlsx": load_csv,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
}


# -----------------------------
# Enhanced Vector Store with Chroma
# -----------------------------

class EnhancedVectorStore:
    """Enhanced vector store using Chroma with LangChain integration."""
    
    def __init__(self, collection_name: str = "documents", embedding_model: str = "sentence-transformers/all-MiniLM-L6-v2", use_openai_embeddings: bool = False):
        self.collection_name = collection_name
        self.embedding_model = embedding_model
        self.use_openai_embeddings = use_openai_embeddings
        self.vectorstore = None
        self.embeddings = None
        self._initialize_embeddings()
        self._initialize_vectorstore()
    
    def _initialize_embeddings(self):
        """Initialize embedding model."""
        if self.use_openai_embeddings and os.getenv("OPENAI_API_KEY"):
            try:
                self.embeddings = OpenAIEmbeddings()
            except Exception as e:
                st.warning(f"Failed to initialize OpenAI embeddings: {e}. Falling back to SentenceTransformer.")
                self.embeddings = SentenceTransformerEmbeddings(model_name=self.embedding_model)
        else:
            self.embeddings = SentenceTransformerEmbeddings(model_name=self.embedding_model)
    
    def _initialize_vectorstore(self):
        """Initialize Chroma vector store."""
        try:
            # Create persistent Chroma client
            chroma_client = chromadb.PersistentClient(path="./chroma_db")
            
            self.vectorstore = Chroma(
                client=chroma_client,
                collection_name=self.collection_name,
                embedding_function=self.embeddings,
            )
        except Exception as e:
            st.error(f"Failed to initialize Chroma vector store: {e}")
            st.stop()
    
    def add_documents(self, documents: List[Document]):
        """Add documents to the vector store."""
        if not documents:
            return
        
        try:
            self.vectorstore.add_documents(documents)
        except Exception as e:
            st.error(f"Failed to add documents to vector store: {e}")
    
    def similarity_search_with_score(self, query: str, k: int = 5):
        """Search for similar documents with scores."""
        try:
            return self.vectorstore.similarity_search_with_score(query, k=k)
        except Exception as e:
            st.error(f"Failed to search vector store: {e}")
            return []
    
    def get_retriever(self, search_kwargs: Dict = None):
        """Get a retriever for use with LangChain chains."""
        if search_kwargs is None:
            search_kwargs = {"k": 5}
        return self.vectorstore.as_retriever(search_kwargs=search_kwargs)
    
    def get_collection_count(self) -> int:
        """Get the number of documents in the collection."""
        try:
            return self.vectorstore._collection.count()
        except:
            return 0


# -----------------------------
# Enhanced Q&A System with LangChain
# -----------------------------

class ConversationalQASystem:
    """Enhanced Q&A system with conversation memory and better source citations."""
    
    def __init__(self, vector_store: EnhancedVectorStore, memory_window: int = 10):
        self.vector_store = vector_store
        self.memory_window = memory_window
        self.conversation_memory = ConversationBufferWindowMemory(
            k=memory_window,
            memory_key="chat_history",
            return_messages=True,
            output_key="answer"
        )
        self.qa_chain = None
        self._initialize_qa_chain()
    
    def _initialize_qa_chain(self):
        """Initialize the conversational QA chain."""
        api_key = os.getenv("OPENAI_API_KEY")
        
        if api_key:
            try:
                # Use OpenAI for generation
                llm = ChatOpenAI(
                    model_name=os.getenv("OPENAI_MODEL", "gpt-4o-mini"),
                    temperature=0.2,
                    api_key=api_key
                )
                
                # Custom prompt template with better source citations
                custom_prompt = PromptTemplate(
                    template="""You are an intelligent document assistant. Use the provided context to answer the user's question accurately and concisely.

IMPORTANT INSTRUCTIONS:
1. Always cite your sources using the format [Source: filename] when referencing information
2. If information comes from multiple sources, cite all relevant sources
3. If you cannot find the answer in the provided context, clearly state that you don't have enough information
4. Provide specific, actionable answers when possible
5. Consider the conversation history for context

Chat History:
{chat_history}

Context from documents:
{context}

Human Question: {question}
""".strip(),
                    input_variables=["chat_history", "context", "question"]
                )
                
                # Create the conversational retrieval chain
                self.qa_chain = ConversationalRetrievalChain.from_llm(
                    llm=llm,
                    retriever=self.vector_index.vectorstore.as_retriever(search_kwargs={"k": 5}),
                    combine_docs_chain_kwargs={"prompt": custom_prompt},
                    return_source_documents=True,
                    verbose=True
                )
                
            except Exception as e:
                st.error(f"Failed to initialize OpenAI: {e}")
                self.qa_chain = None
        else:
            self.qa_chain = None

    def get_answer(self, question: str, chat_history: List[Tuple[str, str]] = None) -> Dict:
        """Get answer using the QA chain."""
        if not self.qa_chain:
            return {
                "answer": "Please provide an OpenAI API key to enable answer generation.",
                "source_documents": []
            }
        
        if chat_history is None:
            chat_history = []
        
        try:
            result = self.qa_chain({
                "question": question,
                "chat_history": chat_history
            })
            return result
        except Exception as e:
            return {
                "answer": f"Error generating answer: {str(e)}",
                "source_documents": []
            }