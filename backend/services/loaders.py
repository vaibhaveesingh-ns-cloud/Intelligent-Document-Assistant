"""Utilities for converting uploaded files into plain text."""

from __future__ import annotations

import io
from pathlib import Path
from typing import Callable, Dict

import pandas as pd
from PIL import Image

try:
    import pdfplumber
except Exception:  # pragma: no cover - optional dependency fallback
    pdfplumber = None

try:
    import docx  # type: ignore
except Exception:  # pragma: no cover
    docx = None

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None


class UnsupportedFileTypeError(RuntimeError):
    """Raised when the requested file type cannot be processed."""


class MissingDependencyError(RuntimeError):
    """Raised when a loader dependency is not available."""


def _require(condition: bool, message: str) -> None:
    if not condition:
        raise MissingDependencyError(message)


def _load_pdf(file_bytes: bytes) -> str:
    _require(pdfplumber is not None, "pdfplumber is required for PDF ingestion.")
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text_parts.append(page.extract_text() or "")
    return "\n".join(text_parts)


def _load_docx(file_bytes: bytes) -> str:
    _require(docx is not None, "python-docx is required for DOCX ingestion.")
    document = docx.Document(io.BytesIO(file_bytes))  # type: ignore[attr-defined]
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _load_pptx(file_bytes: bytes) -> str:
    _require(Presentation is not None, "python-pptx is required for PPTX ingestion.")
    presentation = Presentation(io.BytesIO(file_bytes))
    slides = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slides.append(shape.text)
    return "\n".join(slides)


def _load_image(file_bytes: bytes) -> str:
    _require(
        pytesseract is not None,
        "pytesseract and the Tesseract OCR binary are required for image ingestion.",
    )
    with Image.open(io.BytesIO(file_bytes)) as image:
        return pytesseract.image_to_string(image)


def _load_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except Exception:  # pragma: no cover - fallback path
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _load_md(file_bytes: bytes) -> str:
    return _load_txt(file_bytes)


def _load_csv(file_bytes: bytes) -> str:
    buffer = io.BytesIO(file_bytes)
    try:
        df = pd.read_csv(buffer)
    except Exception:
        buffer.seek(0)
        df = pd.read_excel(buffer)
    return df.to_csv(index=False)


_LOADER_MAP: Dict[str, Callable[[bytes], str]] = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".txt": _load_txt,
    ".md": _load_md,
    ".csv": _load_csv,
    ".png": _load_image,
    ".jpg": _load_image,
    ".jpeg": _load_image,
}


def extract_text_from_bytes(filename: str, file_bytes: bytes) -> str:
    """Return normalized text extracted from the uploaded file."""

    extension = Path(filename).suffix.lower()
    if extension not in _LOADER_MAP:
        raise UnsupportedFileTypeError(
            f"Unsupported file extension '{extension}'. "
            "Supported types: PDF, DOCX, PPTX, TXT, MD, CSV, PNG, JPG, JPEG."
        )

    loader = _LOADER_MAP[extension]
    text = loader(file_bytes)
    return text.strip()


SUPPORTED_EXTENSIONS = tuple(sorted(_LOADER_MAP.keys()))
